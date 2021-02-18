import os
from docx import Document
from docx.shared import Pt, Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


DEMO_FILENAME = 'demo.docx'


if os.path.exists(DEMO_FILENAME):
    print 'Removing old document...'
    os.remove(DEMO_FILENAME)

print 'Generating document...'
document = Document()
document.add_heading('Charts in Word', 0)
p = document.add_paragraph('An example chart:')

chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
x, y, cx, cy = Cm(2), Cm(2), Cm(15), Cm(10)
chart = document.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)

chart.chart_style = 3  # Makes the chart monochrome
chart.has_title = True
chart_title = chart.chart_title
text_frame = chart_title.text_frame
text_frame.text = 'Chart Title Heres'
paragraphs = text_frame.paragraphs
paragraph = paragraphs[0]
paragraph.font.size = Pt(18)
chart.value_axis.has_major_gridlines = False
chart.value_axis.tick_labels.font.size = Pt(9)
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(11)

document.save(DEMO_FILENAME)
print 'Done'
